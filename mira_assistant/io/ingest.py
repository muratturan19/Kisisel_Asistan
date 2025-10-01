"""Document ingestion pipeline for offline processing."""
from __future__ import annotations

import datetime as dt
import hashlib
import logging
import shutil
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, List, Optional, Sequence

import numpy as np

from config import settings
from mira_assistant.core.summarizer import generate_summary
from mira_assistant.core.vector_store import VectorStore
from mira_assistant.core.storage import (
    Chunk,
    Document,
    get_document_by_checksum,
    get_session,
)

LOGGER = logging.getLogger(__name__)


@dataclass
class IngestResult:
    document: Optional[Document]
    chunk_texts: List[str]
    summary: str
    skipped: bool = False


class DocumentIngestor:
    """Process files from the inbox directory into the knowledge base."""

    def __init__(self, *, vector_store: Optional[VectorStore] = None, model_name: Optional[str] = None) -> None:
        self.vector_store = vector_store or VectorStore()
        self._model_name = model_name or settings.embed_model_name
        self._model = None

    def process_inbox(self, default_topic: Optional[str] = None) -> List[str]:
        processed: List[str] = []
        for path in settings.inbox_path.iterdir():
            if not path.is_file():
                continue
            try:
                result = self.ingest(path, topic=default_topic)
            except Exception as exc:  # pragma: no cover - ingestion robustness
                LOGGER.error("Failed to ingest %s: %s", path, exc)
                continue
            if result and not result.skipped and result.document is not None:
                processed.append(result.document.title)
        return processed

    def ingest(self, source_path: Path, topic: Optional[str] = None) -> IngestResult:
        source_path = source_path.expanduser().resolve()
        if not source_path.exists():
            raise FileNotFoundError(source_path)
        checksum = _sha256_of_path(source_path)
        with get_session() as session:
            existing = get_document_by_checksum(session, checksum)
            if existing is not None:
                LOGGER.info("Skipping %s; already ingested", source_path)
                if hasattr(existing, "model_dump"):
                    snapshot = Document.model_validate(existing, from_attributes=True)  # type: ignore[attr-defined]
                else:
                    snapshot = existing
                return IngestResult(document=snapshot, chunk_texts=[], summary="", skipped=True)
        topic = topic or infer_topic_from_filename(source_path.name)
        archive_path = build_archive_path(source_path, topic)
        archive_path.parent.mkdir(parents=True, exist_ok=True)
        shutil.move(str(source_path), archive_path)
        text = extract_text(archive_path)
        chunk_texts = create_chunks(text)
        embeddings = self._embed(chunk_texts)
        tags = generate_tags(text)
        summary = generate_summary(topic, chunk_texts)
        with get_session() as session:
            document = Document(
                path=str(archive_path),
                title=archive_path.stem,
                topic=topic,
                tags={"auto": tags},
                checksum=checksum,
                text_chars=len(text),
                lang="tr",
                ingested_at=dt.datetime.now(dt.timezone.utc),
            )
            session.add(document)
            session.commit()
            session.refresh(document)
            chunks = _create_chunk_models(document.id, chunk_texts, embeddings)
            for chunk in chunks:
                session.add(chunk)
            session.commit()
            if hasattr(Document, "model_validate"):
                document_snapshot = Document.model_validate(document, from_attributes=True)  # type: ignore[attr-defined]
            else:  # pragma: no cover - legacy pydantic v1 fallback
                document_snapshot = Document.from_orm(document)  # type: ignore[attr-defined]
        document_id = document_snapshot.id
        metadata = [{"doc_id": str(document_id), "topic": topic} for _ in chunk_texts]
        ids = [f"doc-{document_id}-chunk-{idx}" for idx in range(len(chunk_texts))]
        self.vector_store.add_embeddings(embeddings.tolist(), metadatas=metadata, ids=ids, documents=chunk_texts)
        return IngestResult(document=document_snapshot, chunk_texts=chunk_texts, summary=summary)

    def _load_model(self):  # type: ignore[no-untyped-def]
        if self._model is None:
            from sentence_transformers import SentenceTransformer  # type: ignore

            self._model = SentenceTransformer(self._model_name)
        return self._model

    def _embed(self, texts: Sequence[str]) -> np.ndarray:
        if not texts:
            return np.zeros((0, 384), dtype=np.float32)
        if settings.offline_only:
            return self._offline_embed(texts)
        model = self._load_model()
        embeddings = model.encode(list(texts), show_progress_bar=False, normalize_embeddings=True)
        return np.asarray(embeddings, dtype=np.float32)

    def _offline_embed(self, texts: Sequence[str]) -> np.ndarray:
        """Generate deterministic embeddings without external dependencies."""

        dim = 384
        vectors = np.zeros((len(texts), dim), dtype=np.float32)
        for row, text in enumerate(texts):
            if not text:
                continue
            tokens = text.lower().split()
            for token in tokens:
                digest = hashlib.sha256(token.encode("utf-8")).digest()
                index = int.from_bytes(digest[:4], "little") % dim
                vectors[row, index] += 1.0
            norm = np.linalg.norm(vectors[row])
            if norm:
                vectors[row] /= norm
        return vectors


def infer_topic_from_filename(filename: str) -> str:
    stem = Path(filename).stem
    parts = stem.split("_")
    if parts:
        return parts[0].capitalize()
    return "Genel"


def build_archive_path(source: Path, topic: str) -> Path:
    now = dt.datetime.now()
    archive_dir = settings.archive_root / topic / now.strftime("%Y-%m")
    return archive_dir / source.name


def _sha256_of_path(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(8192), b""):
            digest.update(chunk)
    return digest.hexdigest()


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(path)
    if suffix in {".docx"}:
        return _extract_docx(path)
    if suffix in {".ppt", ".pptx"}:
        return _extract_pptx(path)
    if suffix in {".png", ".jpg", ".jpeg", ".tiff"}:
        return _extract_image(path)
    return path.read_text(encoding="utf-8", errors="ignore")


def create_chunks(text: str, *, min_tokens: int = 800, max_tokens: int = 1200) -> List[str]:
    words = text.split()
    if not words:
        return [""]
    chunk_size = max(min_tokens, min(max_tokens, 900))
    overlap = 100
    chunks: List[str] = []
    start = 0
    while start < len(words):
        end = min(len(words), start + chunk_size)
        chunk_words = words[start:end]
        chunks.append(" ".join(chunk_words))
        if end == len(words):
            break
        start = end - overlap
    return chunks


def _create_chunk_models(doc_id: int, texts: Sequence[str], embeddings: np.ndarray) -> List[Chunk]:
    chunk_models: List[Chunk] = []
    for idx, text in enumerate(texts):
        embedding = embeddings[idx] if idx < len(embeddings) else np.zeros(384, dtype=np.float32)
        chunk_models.append(
            Chunk(
                doc_id=doc_id,
                seq=idx,
                text=text,
                embedding=embedding.tobytes(),
                tokens=len(text.split()),
            )
        )
    return chunk_models


def _extract_pdf(path: Path) -> str:
    try:
        import pypdf

        reader = pypdf.PdfReader(str(path))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception as exc:  # pragma: no cover - best effort
        LOGGER.warning("PDF parse failed for %s: %s", path, exc)
        return ""


def _extract_docx(path: Path) -> str:
    try:
        import docx

        document = docx.Document(str(path))
        return "\n".join(paragraph.text for paragraph in document.paragraphs)
    except Exception as exc:  # pragma: no cover - best effort
        LOGGER.warning("DOCX parse failed for %s: %s", path, exc)
        return ""


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation

        presentation = Presentation(str(path))
        lines: List[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    lines.append(shape.text)
        return "\n".join(lines)
    except Exception as exc:  # pragma: no cover - best effort
        LOGGER.warning("PPTX parse failed for %s: %s", path, exc)
        return ""


def _extract_image(path: Path) -> str:
    try:
        from PIL import Image
        import pytesseract

        with Image.open(path) as image:
            return pytesseract.image_to_string(image, lang="tur")
    except Exception as exc:  # pragma: no cover - OCR optional
        LOGGER.warning("OCR failed for %s: %s", path, exc)
        return ""


def generate_tags(text: str, limit: int = 5) -> List[str]:
    words = [word.lower() for word in text.split() if len(word) > 4]
    freq: dict[str, int] = {}
    for word in words:
        freq[word] = freq.get(word, 0) + 1
    sorted_items = sorted(freq.items(), key=lambda item: item[1], reverse=True)
    return [word for word, _ in sorted_items[:limit]]


__all__ = [
    "DocumentIngestor",
    "IngestResult",
    "infer_topic_from_filename",
    "build_archive_path",
]
